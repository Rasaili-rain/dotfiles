#!/usr/bin/env python3
"""
Convert PPTX, PPT, and TXT files to PDF format.
Usage: python convert_to_pdf.py <folder_path>
"""

import sys
import os
from pathlib import Path
from typing import List

try:
    from pptx import Presentation
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import ImageReader
    from PIL import Image
    import io
except ImportError:
    print("Installing required packages...")
    os.system("pip install python-pptx reportlab pillow")
    from pptx import Presentation
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import ImageReader
    from PIL import Image
    import io


def convert_txt_to_pdf(txt_path: Path, pdf_path: Path) -> bool:
    """Convert text file to PDF with better formatting."""
    try:
        doc = SimpleDocTemplate(str(pdf_path), pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        with open(txt_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        # Split into paragraphs
        paragraphs = content.split('\n\n')
        
        for para in paragraphs:
            if para.strip():
                p = Paragraph(para.replace('\n', '<br/>'), styles['Normal'])
                story.append(p)
                story.append(Spacer(1, 0.2*inch))
        
        doc.build(story)
        return True
    except Exception as e:
        print(f"Error converting {txt_path}: {e}")
        return False


def convert_markdown_to_pdf(md_path: Path, pdf_path: Path) -> bool:
    """Convert Markdown file to PDF."""
    try:
        import markdown
        from reportlab.lib.styles import ParagraphStyle
        from reportlab.lib.enums import TA_LEFT
        
        doc = SimpleDocTemplate(str(pdf_path), pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        with open(md_path, 'r', encoding='utf-8', errors='ignore') as f:
            md_content = f.read()
        
        # Convert markdown to HTML
        html_content = markdown.markdown(md_content)
        
        # Simple parsing (basic support)
        lines = html_content.split('\n')
        for line in lines:
            if line.strip():
                p = Paragraph(line, styles['Normal'])
                story.append(p)
                story.append(Spacer(1, 0.1*inch))
        
        doc.build(story)
        return True
    except ImportError:
        print("Markdown support requires 'markdown' package. Falling back to text conversion...")
        return convert_txt_to_pdf(md_path, pdf_path)
    except Exception as e:
        print(f"Error converting {md_path}: {e}")
        return False


def convert_with_libreoffice(input_path: Path, output_dir: Path) -> bool:
    """Convert documents using LibreOffice (for DOC, DOCX, ODT, ODP, old PPT)."""
    try:
        # Check if LibreOffice is available
        result = subprocess.run(
            ['libreoffice', '--version'],
            capture_output=True,
            timeout=5
        )
        
        if result.returncode != 0:
            return False
        
        # Convert using LibreOffice headless mode
        cmd = [
            'libreoffice',
            '--headless',
            '--convert-to', 'pdf',
            '--outdir', str(output_dir),
            str(input_path)
        ]
        
        result = subprocess.run(
            cmd,
            capture_output=True,
            timeout=60,
            text=True
        )
        
        return result.returncode == 0
        
    except (subprocess.TimeoutExpired, FileNotFoundError):
        return False
    except Exception as e:
        print(f"LibreOffice conversion error: {e}")
        return False


def convert_pptx_to_pdf(pptx_path: Path, pdf_path: Path) -> bool:
    """Convert PPTX file to PDF."""
    try:
        prs = Presentation(str(pptx_path))
        c = canvas.Canvas(str(pdf_path), pagesize=A4)
        width, height = A4
        
        for slide_num, slide in enumerate(prs.slides):
            if slide_num > 0:
                c.showPage()
            
            # Add slide title or number
            c.setFont("Helvetica-Bold", 16)
            c.drawString(50, height - 50, f"Slide {slide_num + 1}")
            
            # Extract text from shapes
            y_position = height - 100
            c.setFont("Helvetica", 12)
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_lines = shape.text.split('\n')
                    for line in text_lines:
                        if y_position < 50:
                            break
                        c.drawString(50, y_position, line[:80])
                        y_position -= 20
        
        c.save()
        return True
    except Exception as e:
        print(f"Error converting {pptx_path}: {e}")
        return False



def main():
    if len(sys.argv) < 2:
        print("Usage: python convert_to_pdf.py <source_folder> [destination_folder]")
        sys.exit(1)
    
    source_folder = Path(sys.argv[1])
    
    # Determine destination folder
    if len(sys.argv) >= 3:
        dest_folder = Path(sys.argv[2])
        dest_folder.mkdir(parents=True, exist_ok=True)
    else:
        dest_folder = None  # Convert in place
    
    if not source_folder.exists():
        print(f"Error: Folder '{source_folder}' does not exist")
        sys.exit(1)
    
    if not source_folder.is_dir():
        print(f"Error: '{source_folder}' is not a directory")
        sys.exit(1)
    
    # Define supported file extensions
    presentation_exts = ['.pptx', '.ppt', '.odp']
    document_exts = ['.doc', '.docx', '.odt', '.rtf']
    text_exts = ['.txt', '.md']
    
    all_extensions = presentation_exts + document_exts + text_exts
    
    # Find all relevant files
    print(f"Scanning folder: {source_folder}")
    all_files = find_files(source_folder, all_extensions)
    
    print(f"Found {len(all_files)} file(s) to convert")
    
    if len(all_files) == 0:
        print("No files to convert")
        return
    
    if dest_folder:
        print(f"Output folder: {dest_folder}")
    else:
        print("Converting files in place")
    
    converted = 0
    failed = 0
    
    # Group files by type for better processing
    for file_path in all_files:
        ext = file_path.suffix.lower()
        
        # Determine output path - ALWAYS preserve relative location
        if dest_folder:
            # Preserve full directory structure relative to source
            try:
                rel_path = file_path.relative_to(source_folder)
                output_dir = dest_folder / rel_path.parent
                output_dir.mkdir(parents=True, exist_ok=True)
                pdf_path = output_dir / f"{file_path.stem}.pdf"
            except ValueError:
                # Fallback if relative_to fails
                pdf_path = dest_folder / f"{file_path.stem}.pdf"
        else:
            # Convert in place - PDF goes in same directory as source file
            pdf_path = file_path.parent / f"{file_path.stem}.pdf"
        
        # Show relative path for clarity
        if dest_folder:
            display_path = f"{file_path.relative_to(source_folder)} -> {pdf_path.relative_to(dest_folder)}"
        else:
            display_path = f"{file_path.relative_to(source_folder)} -> {pdf_path.relative_to(source_folder)}"
        
        print(f"Converting: {display_path}")
        
        success = False
        
        # Try appropriate conversion method
        if ext == '.pptx':
            success = convert_pptx_to_pdf(file_path, pdf_path)
        elif ext == '.txt':
            success = convert_txt_to_pdf(file_path, pdf_path)
        elif ext == '.md':
            success = convert_markdown_to_pdf(file_path, pdf_path)
        elif ext in ['.ppt', '.odp', '.doc', '.docx', '.odt', '.rtf']:
            # Try LibreOffice first
            output_dir = pdf_path.parent
            success = convert_with_libreoffice(file_path, output_dir)
            
            if not success:
                print(f"  ⚠ LibreOffice not available for {ext} files")
                print(f"  Install LibreOffice: sudo apt install libreoffice")
        
        if success:
            converted += 1
            print(f"  ✓ Success")
        else:
            failed += 1
            print(f"  ✗ Failed")
    
    print(f"\n{'='*50}")
    print(f"Conversion complete!")
    print(f"Successfully converted: {converted}")
    print(f"Failed: {failed}")
    print(f"{'='*50}")


if __name__ == "__main__":
    main()