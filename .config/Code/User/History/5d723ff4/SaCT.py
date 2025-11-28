#!/usr/bin/env python3
"""
Convert PPTX, PPT, and TXT files to PDF format.
Usage: python convert_to_pdf.py <folder_path>
"""

import sys
import os
from pathlib import Path
from typing import List
from pptx import Presentation
from reportlab.lib.pagesizes import letter, A4
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
from PIL import Image
import io



def convert_txt_to_pdf(txt_path: Path, pdf_path: Path) -> bool:
    """Convert text file to PDF."""
    try:
        c = canvas.Canvas(str(pdf_path), pagesize=letter)
        width, height = letter
        
        with open(txt_path, 'r', encoding='utf-8', errors='ignore') as f:
            lines = f.readlines()
        
        y = height - 50
        for line in lines:
            line = line.rstrip('\n')
            if y < 50:
                c.showPage()
                y = height - 50
            
            c.drawString(50, y, line[:100])  # Limit line length
            y -= 15
        
        c.save()
        return True
    except Exception as e:
        print(f"Error converting {txt_path}: {e}")
        return False


def convert_pptx_to_pdf(pptx_path: Path, pdf_path: Path) -> bool:
    """Convert PPTX file to PDF."""
    try:
        prs = Presentation(str(pptx_path))
        c = canvas.Canvas(str(pdf_path), pagesize=A4)
        width, height = A4
        
        for slide_num, slide in enumerate(prs.slides):
            if slide_num > 0:
                c.showPage()
            
            # Add slide title or number
            c.setFont("Helvetica-Bold", 16)
            c.drawString(50, height - 50, f"Slide {slide_num + 1}")
            
            # Extract text from shapes
            y_position = height - 100
            c.setFont("Helvetica", 12)
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_lines = shape.text.split('\n')
                    for line in text_lines:
                        if y_position < 50:
                            break
                        c.drawString(50, y_position, line[:80])
                        y_position -= 20
        
        c.save()
        return True
    except Exception as e:
        print(f"Error converting {pptx_path}: {e}")
        return False


def find_files(folder_path: Path, extensions: List[str]) -> List[Path]:
    """Recursively find all files with given extensions."""
    files = []
    for ext in extensions:
        files.extend(folder_path.rglob(f"*{ext}"))
    return files


def main():
    if len(sys.argv) < 2:
        print("Usage: python convert_to_pdf.py <folder_path>")
        sys.exit(1)
    
    folder_path = Path(sys.argv[1])
    
    if not folder_path.exists():
        print(f"Error: Folder '{folder_path}' does not exist")
        sys.exit(1)
    
    if not folder_path.is_dir():
        print(f"Error: '{folder_path}' is not a directory")
        sys.exit(1)
    
    # Find all relevant files
    print(f"Scanning folder: {folder_path}")
    slide_extensions = ['.pptx', '.ppt']
    txt_extensions = ['.txt']
    
    slide_files = find_files(folder_path, slide_extensions)
    txt_files = find_files(folder_path, txt_extensions)
    
    total_files = len(slide_files) + len(txt_files)
    print(f"Found {len(slide_files)} slide file(s) and {len(txt_files)} text file(s)")
    
    if total_files == 0:
        print("No files to convert")
        return
    
    converted = 0
    failed = 0
    
    # Convert slide files
    for slide_file in slide_files:
        pdf_path = slide_file.with_suffix('.pdf')
        print(f"Converting: {slide_file.name} -> {pdf_path.name}")
        
        if convert_pptx_to_pdf(slide_file, pdf_path):
            converted += 1
        else:
            failed += 1
    
    # Convert text files
    for txt_file in txt_files:
        pdf_path = txt_file.with_suffix('.pdf')
        print(f"Converting: {txt_file.name} -> {pdf_path.name}")
        
        if convert_txt_to_pdf(txt_file, pdf_path):
            converted += 1
        else:
            failed += 1
    
    print(f"\nConversion complete!")
    print(f"Successfully converted: {converted}")
    print(f"Failed: {failed}")


if __name__ == "__main__":
    main()